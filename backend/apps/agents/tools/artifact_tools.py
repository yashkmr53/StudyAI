"""Artifact generation tools (Phase 4).

Provides tools for generating PDF, PPTX, and DOCX artifacts from
agentic workflows.
"""
import logging
import os
import tempfile
from typing import Any

from apps.agents.tools.base import ToolMetadata, ToolInput, ToolOutput
from apps.agents.tools.learning import BaseStudyAITool

logger = logging.getLogger(__name__)


class GeneratePDFInput(ToolInput):
    title: str
    content: str
    sections: list[dict] = []


class GeneratePPTXInput(ToolInput):
    title: str
    slides: list[dict]


class GenerateDOCXInput(ToolInput):
    title: str
    content: str
    sections: list[dict] = []


class GeneratePDFTool(BaseStudyAITool):
    metadata = ToolMetadata(
        name="generate_pdf",
        description="Generate a PDF document from structured content. Use for revision materials, summaries, or study guides.",
        category="artifact",
        input_schema=GeneratePDFInput,
        output_schema=ToolOutput,
        requires_auth=True,
        timeout_seconds=60,
    )

    def _execute(self, input: GeneratePDFInput, *, user, request_id: str) -> ToolOutput:
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
        except ImportError:
            return ToolOutput(
                success=False,
                result={},
                error="PDF generation requires reportlab. Install with: pip install reportlab",
            )

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            doc = SimpleDocTemplate(tmp.name, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            story.append(Paragraph(input.title, styles["Title"]))
            story.append(Spacer(1, 12))
            story.append(Paragraph(input.content, styles["BodyText"]))
            for section in input.sections:
                story.append(Spacer(1, 12))
                story.append(Paragraph(section.get("heading", ""), styles["Heading2"]))
                story.append(Paragraph(section.get("body", ""), styles["BodyText"]))

            doc.build(story)
            tmp_path = tmp.name

        try:
            from django.core.files.base import ContentFile
            from apps.documents.models import Document

            with open(tmp_path, "rb") as f:
                data = f.read()

            document = Document.objects.create(
                profile=user.profile,
                filename=f"{input.title.replace(' ', '_')}.pdf",
                source_type="pdf",
                file_size=len(data),
            )

            from apps.storage.services import StorageService
            storage = StorageService()
            storage.save(f"documents/{document.pk}/original.pdf", ContentFile(data))

            os.unlink(tmp_path)

            return ToolOutput(
                success=True,
                result={
                    "document_id": str(document.pk),
                    "filename": document.filename,
                    "url": f"/api/v1/documents/{document.pk}",
                },
            )
        except Exception as exc:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
            return ToolOutput(success=False, result={}, error=str(exc))


class GeneratePPTXTool(BaseStudyAITool):
    metadata = ToolMetadata(
        name="generate_pptx",
        description="Generate a PowerPoint presentation from structured slides. Use for presentations about study topics.",
        category="artifact",
        input_schema=GeneratePPTXInput,
        output_schema=ToolOutput,
        requires_auth=True,
        timeout_seconds=60,
    )

    def _execute(self, input: GeneratePPTXInput, *, user, request_id: str) -> ToolOutput:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return ToolOutput(
                success=False,
                result={},
                error="PPTX generation requires python-pptx. Install with: pip install python-pptx",
            )

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = input.title

            for slide_data in input.slides:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_data.get("title", "")
                slide.placeholders[1].text = slide_data.get("content", "")

            prs.save(tmp.name)
            tmp_path = tmp.name

        try:
            from django.core.files.base import ContentFile
            from apps.documents.models import Document

            with open(tmp_path, "rb") as f:
                data = f.read()

            document = Document.objects.create(
                profile=user.profile,
                filename=f"{input.title.replace(' ', '_')}.pptx",
                source_type="pptx",
                file_size=len(data),
            )

            from apps.storage.services import StorageService
            storage = StorageService()
            storage.save(f"documents/{document.pk}/original.pptx", ContentFile(data))

            os.unlink(tmp_path)

            return ToolOutput(
                success=True,
                result={
                    "document_id": str(document.pk),
                    "filename": document.filename,
                    "url": f"/api/v1/documents/{document.pk}",
                },
            )
        except Exception as exc:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
            return ToolOutput(success=False, result={}, error=str(exc))


class GenerateDOCXTool(BaseStudyAITool):
    metadata = ToolMetadata(
        name="generate_docx",
        description="Generate a Word document from structured content. Use for detailed notes, essays, or reports.",
        category="artifact",
        input_schema=GenerateDOCXInput,
        output_schema=ToolOutput,
        requires_auth=True,
        timeout_seconds=60,
    )

    def _execute(self, input: GenerateDOCXInput, *, user, request_id: str) -> ToolOutput:
        try:
            from docx import Document as DocxDocument
            from docx.shared import Inches
        except ImportError:
            return ToolOutput(
                success=False,
                result={},
                error="DOCX generation requires python-docx. Install with: pip install python-docx",
            )

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            doc = DocxDocument()
            doc.add_heading(input.title, 0)
            doc.add_paragraph(input.content)
            for section in input.sections:
                doc.add_heading(section.get("heading", ""), level=1)
                doc.add_paragraph(section.get("body", ""))
            doc.save(tmp.name)
            tmp_path = tmp.name

        try:
            from django.core.files.base import ContentFile
            from apps.documents.models import Document

            with open(tmp_path, "rb") as f:
                data = f.read()

            document = Document.objects.create(
                profile=user.profile,
                filename=f"{input.title.replace(' ', '_')}.docx",
                source_type="docx",
                file_size=len(data),
            )

            from apps.storage.services import StorageService
            storage = StorageService()
            storage.save(f"documents/{document.pk}/original.docx", ContentFile(data))

            os.unlink(tmp_path)

            return ToolOutput(
                success=True,
                result={
                    "document_id": str(document.pk),
                    "filename": document.filename,
                    "url": f"/api/v1/documents/{document.pk}",
                },
            )
        except Exception as exc:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
            return ToolOutput(success=False, result={}, error=str(exc))
